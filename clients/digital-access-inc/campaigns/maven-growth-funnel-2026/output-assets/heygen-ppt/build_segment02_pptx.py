"""Build HeyGen PPTX Segment 2. Run: python build_segment02_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide
from build_segment01_pptx import load_notes_from_file

OUT = Path(__file__).parent / "Maven-HS-Seg02-Gaps-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-02-gaps.txt"

SLIDES = [
    {"title": "The 5 Revenue Gaps", "subtitle": "Bleeding money from local service businesses"},
    {"title": "Gap #1: The Unanswered Call", "subtitle": ""},
    {"title": "$468,000 per year", "subtitle": "Lost to missed calls", "accent_word": "$468,000"},
    {"title": "35–40% of calls unanswered", "subtitle": "During regular business hours", "accent_word": "35–40%"},
    {"title": "Gap #2: The Review Gap", "subtitle": "Google “near me” rankings & trust"},
    {"title": "244-review gap", "subtitle": "HVAC Phoenix example", "accent_word": "244"},
    {"title": "Gap #3: The Visibility Problem", "subtitle": "Page 1 captures ~75% of clicks"},
    {"title": "$119,000 invisible", "subtitle": "One keyword — plumbing, Austin TX", "accent_word": "$119,000"},
    {"title": "Gap #4: Speed-to-Lead", "subtitle": "21× higher conversion in 5 min vs 30 min", "accent_word": "21×"},
    {"title": "$230,400 per year", "subtitle": "Cost of slow web-lead response", "accent_word": "$230,400"},
    {"title": "Gap #5: The Data Vacuum", "subtitle": "Your portal — gap scores, keywords, roadmap"},
]


def load_notes():
    return load_notes_from_file(VOICE)


def main():
    notes = load_notes()
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        add_slide(prs, blank, {**spec, "notes": note})
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
