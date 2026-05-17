"""Shared helpers for Maven HeyGen PPTX builders."""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
MUTED = RGBColor(0x64, 0x74, 0x8B)

TEXT_LEFT = Inches(0.75)
TEXT_TOP = Inches(1.85)
TEXT_WIDTH = Inches(8.75)
TEXT_HEIGHT = Inches(4.5)
FOOTER_TOP = Inches(6.85)


def set_dark_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_NAVY


def add_footer(slide):
    box = slide.shapes.add_textbox(TEXT_LEFT, FOOTER_TOP, TEXT_WIDTH, Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = "MAVEN by Digital Access Inc"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.LEFT


def add_slide_text(slide, title, subtitle="", accent_word=None, title_size=None):
    box = slide.shapes.add_textbox(TEXT_LEFT, TEXT_TOP, TEXT_WIDTH, TEXT_HEIGHT)
    tf = box.text_frame
    tf.word_wrap = True
    size = title_size or (Pt(36) if len(title) > 50 else Pt(44))

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.bold = True
    p.font.size = size
    p.font.color.rgb = WHITE

    if accent_word and accent_word.lower() in title.lower():
        idx = title.lower().find(accent_word.lower())
        before = title[:idx]
        word = title[idx : idx + len(accent_word)]
        after = title[idx + len(accent_word) :]
        p.text = before
        if word:
            r = p.add_run()
            r.text = word
            r.font.size = size
            r.font.bold = True
            r.font.color.rgb = EMERALD
        if after:
            r = p.add_run()
            r.text = after
            r.font.size = size
            r.font.bold = True
            r.font.color.rgb = WHITE
    else:
        p.text = title

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(26)
        p2.font.color.rgb = EMERALD
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(16)


def add_slide(prs, blank_layout, spec):
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_slide_text(
        slide,
        spec["title"],
        spec.get("subtitle", ""),
        spec.get("accent_word"),
        spec.get("title_size"),
    )
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = spec["notes"].strip()
    return slide
