"""Build HeyGen PPTX Segment 1. Run: python normalize_tts.py && python build_segment01_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide

OUT = Path(__file__).parent / "Maven-HS-Seg01-Hook-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-01-hook.txt"


def _skip_heygen_scene(header: str) -> bool:
    h = header.upper()
    return "DO NOT USE HEYGEN" in h or "USE CAPCUT" in h


def load_notes_from_file(path: Path):
    """Pull SCENE blocks from a heygen-voice-paste file."""
    text = path.read_text(encoding="utf-8")
    notes = []
    block = []
    skip_scene = False
    for line in text.splitlines():
        if line.startswith("========== SCENE"):
            if block and not skip_scene:
                notes.append("\n".join(block).strip())
            block = []
            skip_scene = _skip_heygen_scene(line)
            continue
        if line.startswith("MAVEN HOME") or line.startswith("HeyGen project"):
            continue
        if " scenes" in line and "paste" not in line.lower():
            continue
        if line.startswith("2 HeyGen"):
            continue
        if skip_scene or not line.strip():
            continue
        block.append(line)
    if block and not skip_scene:
        notes.append("\n".join(block).strip())
    return notes


def load_notes():
    return load_notes_from_file(VOICE)


SLIDES = [
    {"title": "How much more revenue do you want to add to your business this year?", "subtitle": "", "accent_word": "revenue"},
    {"title": "The only question that matters", "subtitle": "Revenue — not clicks"},
    {"title": "In the next 30 minutes", "subtitle": "Your business · Your market · Your numbers"},
    {"title": "Example: Roofing contractor", "subtitle": "$1.2M today → $1.6M goal"},
    {"title": "Maven by Digital Access Inc", "subtitle": "Not a traditional marketing agency", "accent_word": "Maven"},
    {"title": "Today: Real analysis walkthrough", "subtitle": "Plus how to get the same analysis for your business"},
]


def main():
    notes = load_notes()
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        spec = {**spec, "notes": note}
        add_slide(prs, blank, spec)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
