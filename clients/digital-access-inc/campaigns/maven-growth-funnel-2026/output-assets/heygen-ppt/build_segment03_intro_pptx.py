"""Build HeyGen PPTX Segment 3 intro/outro. Run: python build_segment03_intro_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide
from build_segment01_pptx import load_notes_from_file

OUT = Path(__file__).parent / "Maven-HS-Seg03-Demo-Intro-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-03-intro-outro.txt"

SLIDES = [
    {"title": "The 7-Minute Analysis", "subtitle": "Live walkthrough — Fort Worth roofer"},
    {"title": "That's the analysis.", "subtitle": "Yours forever — whether you work with us or not"},
]


def main():
    notes = load_notes_from_file(VOICE)
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        add_slide(prs, blank, {**spec, "notes": note})
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
