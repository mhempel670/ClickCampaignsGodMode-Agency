"""Build HeyGen PPTX Segment 6. Run: python build_segment06_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide
from build_segment01_pptx import load_notes_from_file

OUT = Path(__file__).parent / "Maven-HS-Seg06-Offer-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-06-offer.txt"

SLIDES = [
    {"title": "What to do right now", "subtitle": "Click the button below"},
    {"title": "In 7 minutes you'll have:", "subtitle": "Gap analysis · CAC · Keywords · Roadmap · Portal"},
    {
        "title": "$39",
        "subtitle": "Full customized analysis for your business",
        "accent_word": "$39",
    },
    {
        "title": "After you pay",
        "subtitle": "8 questions → AI builds plan → portal access",
    },
    {
        "title": "Webinar only",
        "subtitle": "$39 now · $297 starting next week",
        "accent_word": "$39",
    },
    {
        "title": "Get Your Custom Maven Analysis",
        "subtitle": "Click Below",
        "accent_word": "Maven",
    },
]


def main():
    notes = load_notes_from_file(VOICE)
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        add_slide(prs, blank, {**spec, "notes": note})
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
