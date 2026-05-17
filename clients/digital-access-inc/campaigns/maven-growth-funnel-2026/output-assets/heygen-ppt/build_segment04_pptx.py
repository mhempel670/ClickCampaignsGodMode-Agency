"""Build HeyGen PPTX Segment 4. Run: python build_segment04_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide
from build_segment01_pptx import load_notes_from_file

OUT = Path(__file__).parent / "Maven-HS-Seg04-Expect-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-04-expect.txt"

SLIDES = [
    {"title": "What happens when you click the button", "subtitle": ""},
    {
        "title": "Steps 1–2",
        "subtitle": "Same interface · 8 questions · Use your real numbers",
    },
    {
        "title": "Step 3",
        "subtitle": "5–7 minute analysis — email when ready",
        "accent_word": "5–7",
    },
    {
        "title": "Steps 4–5",
        "subtitle": "Review plan · Book call · Run it yourself · Or we execute",
    },
]


def main():
    notes = load_notes_from_file(VOICE)
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        add_slide(prs, blank, {**spec, "notes": note})
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
