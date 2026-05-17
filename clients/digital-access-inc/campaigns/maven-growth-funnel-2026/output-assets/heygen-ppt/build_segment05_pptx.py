"""Build HeyGen PPTX Segment 5. Run: python build_segment05_pptx.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from heygen_ppt_common import add_slide
from build_segment01_pptx import load_notes_from_file

OUT = Path(__file__).parent / "Maven-HS-Seg05-Objections-DARK.pptx"
VOICE = Path(__file__).parent.parent / "heygen-voice-paste" / "segment-05-objections.txt"

SLIDES = [
    {"title": "Questions I usually get", "subtitle": ""},
    {"title": '"I already have a marketing agency"', "subtitle": ""},
    {
        "title": '"Can you guarantee results?"',
        "subtitle": "No guarantees — full transparency yes",
    },
    {
        "title": '"Sounds like every other pitch"',
        "subtitle": "$39 full analysis first — then you decide",
        "accent_word": "$39",
    },
    {
        "title": '"I tried SEO before"',
        "subtitle": "Customized to your market — not generic",
    },
    {
        "title": '"I don\'t have time"',
        "subtitle": "The system runs autonomously",
    },
]


def main():
    notes = load_notes_from_file(VOICE)
    if len(notes) != len(SLIDES):
        raise SystemExit(f"Expected {len(SLIDES)} note blocks, got {len(notes)}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for spec, note in zip(SLIDES, notes):
        add_slide(prs, blank, {**spec, "notes": note})
    prs.save(OUT)
    print(f"Wrote {OUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
